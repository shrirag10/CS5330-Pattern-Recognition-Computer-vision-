"""
Revamped final-project deck: clean white theme, flowcharts, results, engaging.
12 slides, 16:9. Georgia headers / Calibri body, crimson accent.
Run: python3 make_pptx.py   ->  presentation.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "presentation.pptx")

# palette
ACCENT   = RGBColor(0x99,0x1B,0x1B)
ACCENT_D = RGBColor(0x7A,0x3A,0x3A)
INK      = RGBColor(0x1E,0x29,0x3B)
MUTED    = RGBColor(0x64,0x74,0x8B)
FAINT    = RGBColor(0x94,0xA3,0xB8)
RULE     = RGBColor(0xE5,0xE7,0xEB)
PANEL    = RGBColor(0xF8,0xFA,0xFC)
SOFT     = RGBColor(0xF7,0xEC,0xEC)
SOFTLINE = RGBColor(0xE3,0xC9,0xC9)
SLATE6   = RGBColor(0x47,0x55,0x69)
SLATE4   = RGBColor(0x94,0xA3,0xB8)
FROZBG   = RGBColor(0xEE,0xF2,0xF7)
FROZLN   = RGBColor(0xD6,0xDE,0xE8)
DARK     = RGBColor(0x0F,0x17,0x2A)
WHITE    = RGBColor(0xFF,0xFF,0xFF)
HOTBG    = RGBColor(0xF3,0xD9,0xD9)

SERIF="Georgia"; SANS="Calibri"

prs=Presentation()
prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]

def slide():
    s=prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb=WHITE
    return s

def _set_radius(shape,frac=0.09):
    try:
        shape.adjustments[0]=frac
    except Exception: pass

def rect(s,l,t,w,h,fill=None,line=None,lw=1.0,rounded=False,radius=0.09,dash=None):
    shp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
                           Inches(l),Inches(t),Inches(w),Inches(h))
    if rounded: _set_radius(shp,radius)
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb=fill
    if line is None: shp.line.fill.background()
    else:
        shp.line.color.rgb=line; shp.line.width=Pt(lw)
        if dash:
            ln=shp.line._get_or_add_ln(); d=ln.makeelement(qn('a:prstDash'),{'val':dash}); ln.append(d)
    shp.shadow.inherit=False
    return shp

def line(s,x1,y1,x2,y2,color=SLATE4,w=1.5):
    c=s.shapes.add_connector(1,Inches(x1),Inches(y1),Inches(x2),Inches(y2))
    c.line.color.rgb=color; c.line.width=Pt(w); c.shadow.inherit=False
    return c

def _run(p,txt,size,color,bold=False,italic=False,font=SANS):
    r=p.add_run(); r.text=txt; f=r.font
    f.size=Pt(size); f.bold=bold; f.italic=italic; f.name=font
    f.color.rgb=color
    return r

def text(s,l,t,w,h,runs,size=16,color=INK,bold=False,italic=False,font=SANS,
         align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,line_sp=1.0,space_after=0):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    for m in ('margin_left','margin_right','margin_top','margin_bottom'): setattr(tf,m,0)
    paras=runs if isinstance(runs,list) else [runs]
    for k,para in enumerate(paras):
        p=tf.paragraphs[0] if k==0 else tf.add_paragraph()
        p.alignment=align; p.line_spacing=line_sp; p.space_before=0
        if space_after: p.space_after=Pt(space_after)
        if isinstance(para,str):
            _run(p,para,size,color,bold,italic,font)
        else:
            for seg in para:
                txt,opts=seg if isinstance(seg,tuple) else (seg,{})
                _run(p,txt,opts.get('size',size),opts.get('color',color),
                     opts.get('bold',bold),opts.get('italic',italic),opts.get('font',font))
    return tb

def bullets(s,l,t,w,h,items,size=15,gap=8,color=SLATE6,bpx=ACCENT):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True
    for m in ('margin_left','margin_right','margin_top','margin_bottom'): setattr(tf,m,0)
    for k,it in enumerate(items):
        p=tf.paragraphs[0] if k==0 else tf.add_paragraph()
        p.line_spacing=1.28; p.space_after=Pt(gap); p.space_before=0
        _run(p,"▪  ",size,bpx,bold=True)
        if isinstance(it,str): _run(p,it,size,color)
        else:
            for seg in it:
                txt,opts=seg if isinstance(seg,tuple) else (seg,{})
                _run(p,txt,opts.get('size',size),opts.get('color',color),opts.get('bold',False),opts.get('italic',False))
    return tb

def kicker(s,txt):
    text(s,0.9,0.62,10,0.35,txt.upper(),size=12.5,color=ACCENT,bold=True)

def h2(s,txt,rule=True):
    text(s,0.9,0.98,11.5,0.9,txt,size=27,color=INK,bold=True,font=SERIF,line_sp=1.05)
    if rule: rect(s,0.92,1.92,0.62,0.045,fill=ACCENT)

def footer(s,left,num):
    text(s,0.9,7.02,8,0.3,left,size=10.5,color=FAINT)
    text(s,12.0,7.0,0.9,0.3,num,size=11,color=FAINT,bold=True,align=PP_ALIGN.RIGHT)

def card(s,l,t,w,h,tag,title,items,hero=False):
    rect(s,l,t,w,h,fill=(SOFT if hero else PANEL),line=(SOFTLINE if hero else RULE),lw=1,rounded=True,radius=0.06)
    text(s,l+0.28,t+0.24,w-0.5,0.3,tag.upper(),size=11.5,color=ACCENT,bold=True)
    text(s,l+0.28,t+0.56,w-0.5,0.4,title,size=17,color=INK,bold=True,font=SERIF)
    if items: bullets(s,l+0.28,t+1.06,w-0.52,h-1.2,items,size=13,gap=6)

def callout(s,l,t,w,h,parts):
    rect(s,l,t,w,h,fill=SOFT,line=SOFTLINE,lw=1,rounded=True,radius=0.05)
    rect(s,l,t,0.07,h,fill=ACCENT)
    text(s,l+0.3,t,w-0.55,h,[parts],size=14.5,color=ACCENT_D,anchor=MSO_ANCHOR.MIDDLE,line_sp=1.25)

def image_centered(s,path,top,width):
    from PIL import Image
    iw,ih=Image.open(path).size
    h=width*ih/iw
    left=(SW-width)/2
    s.shapes.add_picture(path,Inches(left),Inches(top),Inches(width),Inches(h))
    return h

SW=13.333

# ---------------- SLIDE 1 : TITLE ----------------
def dot(s,x,y,d,color):
    o=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(y),Inches(d),Inches(d))
    o.fill.solid(); o.fill.fore_color.rgb=color; o.line.fill.background(); o.shadow.inherit=False
    return o
def chip(s,x,y,w,label,bg,ln,tc,dc):
    rect(s,x,y,w,0.52,fill=bg,line=ln,lw=1.2,rounded=True,radius=0.5)
    dot(s,x+0.22,y+0.195,0.13,dc)
    text(s,x+0.44,y,w-0.5,0.52,label,size=12.5,color=tc,bold=True,anchor=MSO_ANCHOR.MIDDLE)

s=slide()
rect(s,0,0,0.16,7.5,fill=ACCENT)                     # left accent band
text(s,1.0,1.12,10,0.35,"CS 5330 · FINAL PROJECT",size=13,color=ACCENT,bold=True)
text(s,0.98,1.55,11.6,1.9,[
    [("Transfer Learning vs.",{})],
    [("Random Initialization",{'color':ACCENT})],
],size=50,color=INK,bold=True,font=SERIF,line_sp=1.04)
text(s,1.02,3.35,9.9,1.0,
    "“How much of transfer learning’s advantage is the pretrained features — and how much is just fewer knobs to turn?”",
    size=19,color=SLATE6,italic=True,font=SERIF,line_sp=1.3)
# three-condition teaser
chip(s,1.02,4.62,2.55,"Pretrained-Frozen",SOFT,SOFTLINE,ACCENT,ACCENT)
chip(s,3.77,4.62,2.35,"Random-Frozen",FROZBG,FROZLN,SLATE6,SLATE4)
chip(s,6.32,4.62,2.15,"Random-Full",FROZBG,FROZLN,SLATE6,SLATE4)
rect(s,1.0,5.55,11.4,0.02,fill=RULE)
meta=[("PRESENTER","Shriman Raghav Srinivasan"),("DATASET","Intel Image Classification · 6 classes"),("BACKBONE","ResNet-18")]
mx=1.0
for lab,val in meta:
    text(s,mx,5.72,3.95,0.3,lab,size=11,color=FAINT,bold=True)
    text(s,mx,6.04,3.95,0.4,val,size=14.5,color=INK,bold=True)
    mx+=4.02
text(s,1.0,7.02,8,0.3,"Northeastern University · Khoury College",size=11,color=FAINT)
text(s,12.0,7.02,0.9,0.3,"01",size=11,color=FAINT,bold=True,align=PP_ALIGN.RIGHT)

# ---------------- SLIDE 2 : PROBLEM ----------------
s=slide(); kicker(s,"The Problem"); h2(s,"One comparison, two hidden variables")
bullets(s,0.9,2.35,6.0,4.2,[
    [("The usual move: take an ",{}),("ImageNet-pretrained",{'bold':True,'color':INK}),(" network, freeze the early layers, fine-tune the rest.",{})],
    [("It beats training from scratch — but ",{}),("why?",{'bold':True,'color':INK})],
    [("Freezing changes ",{}),("two",{'bold':True,'color':INK}),(" things at once: the ",{}),("quality",{'bold':True,'color':INK}),(" of the frozen features and the ",{}),("number",{'bold':True,'color':INK}),(" of trainable parameters.",{})],
    [("A naive pretrained-vs-scratch test ",{}),("confounds",{'bold':True,'color':INK}),(" them — you can’t tell which one earned the win.",{})],
],size=16,gap=14)
card(s,7.15,2.2,5.25,1.35,"Factor A","Feature quality",None,hero=True)
text(s,7.43,3.02,4.7,0.5,"Pretrained weights encode edges, textures, shapes learned from a million images.",size=13,color=ACCENT_D,line_sp=1.25)
card(s,7.15,3.7,5.25,1.35,"Factor B","Trainable capacity",None)
text(s,7.43,4.52,4.7,0.5,"Freezing early layers also shrinks how many parameters the optimizer can move.",size=13,color=SLATE6,line_sp=1.25)
callout(s,7.15,5.25,5.25,0.8,[("Goal:  ",{'bold':True,'color':ACCENT}),("a design that isolates Factor A from Factor B.",{})])
footer(s,"Transfer Learning vs. Random Init","02")

# ---------------- SLIDE 3 : DESIGN FLOWCHART ----------------
s=slide(); kicker(s,"The Design"); h2(s,"Three conditions, one controlled variable")
conds=[
    (2.55,"Pretrained-Frozen","θf = ImageNet weights, frozen · train θt",SOFT,SOFTLINE,ACCENT,ACCENT_D,None),
    (3.80,"Random-Frozen","θf = random weights, frozen · train θt",FROZBG,FROZLN,SLATE6,MUTED,"CONTROL"),
    (5.05,"Random-Full","everything random · train all layers",FROZBG,FROZLN,RGBColor(0x33,0x41,0x55),MUTED,None),
]
bx=5.4; bw=6.6; bh=1.05
srcx,srcy=3.1,4.325   # right-edge mid of source node
# connectors first so boxes/node sit on top
for ty,*_ in conds:
    line(s,srcx,srcy,bx,ty+bh/2,color=SLATE4,w=1.6)
rect(s,0.95,3.65,2.15,1.35,fill=DARK,line=DARK,rounded=True,radius=0.09)
text(s,0.95,3.83,2.15,0.5,"ResNet-18",size=17,color=WHITE,bold=True,font=SERIF,align=PP_ALIGN.CENTER)
text(s,0.95,4.33,2.15,0.6,"split into early θf + late θt",size=12,color=RGBColor(0xCB,0xD5,0xE1),align=PP_ALIGN.CENTER,line_sp=1.15)
for ty,name,desc,bg,ln,cn,cs,badge in conds:
    rect(s,bx,ty,bw,bh,fill=bg,line=ln,lw=1.5,rounded=True,radius=0.09)
    nm=[(name,{'bold':True,'color':cn,'size':16,'font':SERIF})]
    if badge: nm.append(("    "+badge,{'bold':True,'color':cn,'size':11}))
    text(s,bx+0.3,ty+0.18,bw-0.5,0.4,[nm])
    text(s,bx+0.3,ty+0.6,bw-0.5,0.4,desc,size=12.5,color=cs)
rect(s,0.95,6.28,11.45,0.02,fill=RULE)
text(s,0.95,6.42,5.6,0.5,[[("PT-F ",{'bold':True,'color':ACCENT}),("vs",{'bold':True}),(" R-F  →  isolates ",{}),("feature quality",{'bold':True,'color':INK}),(" (capacity matched)",{})]],size=13.5,color=SLATE6,line_sp=1.2)
text(s,6.9,6.42,5.5,0.5,[[("R-F ",{'bold':True,'color':ACCENT}),("vs",{'bold':True}),(" R-FL  →  isolates ",{}),("trainable capacity",{'bold':True,'color':INK}),(" (both random)",{})]],size=13.5,color=SLATE6,line_sp=1.2)
footer(s,"Transfer Learning vs. Random Init","03")

# ---------------- SLIDE 4 : METHOD SPLIT ----------------
s=slide(); kicker(s,"Method"); h2(s,"Where the network is cut",rule=False)
text(s,0.9,1.9,11.4,0.8,[[("Freeze the generic early layers; train only the task-specific tail. PT-F and R-F share the ",{}),("exact same",{'bold':True,'color':INK}),(" trainable count — that’s what makes them a fair fight.",{})]],size=16,color=MUTED,line_sp=1.35)
blocks=[("conv1","+bn1",True),("layer1","",True),("layer2","",True),("layer3","",True),("layer4","",False),("fc","6-class",False)]
bx=0.9; bw=1.83; gap=0.06; by=2.95; bh=1.15
for name,sub,frozen in blocks:
    rect(s,bx,by,bw,bh,fill=(FROZBG if frozen else SOFT),line=(FROZLN if frozen else SOFTLINE),lw=1.5,rounded=True,radius=0.08)
    text(s,bx,by+0.28,bw,0.4,name,size=15,color=(SLATE6 if frozen else ACCENT),bold=True,font=SERIF,align=PP_ALIGN.CENTER)
    if sub: text(s,bx,by+0.68,bw,0.3,sub,size=11.5,color=(MUTED if frozen else ACCENT_D),align=PP_ALIGN.CENTER)
    bx+=bw+gap
rect(s,0.9,4.32,7.5,0.5,fill=FROZBG,line=FROZLN,lw=1,rounded=True,radius=0.12,dash="dash")
text(s,1.1,4.32,7.2,0.5,[[("θf · frozen   ",{'bold':True,'color':SLATE6,'font':SERIF}),("2,782,784 params · BN forced to eval()",{'color':SLATE6})]],size=13,anchor=MSO_ANCHOR.MIDDLE)
rect(s,8.52,4.32,3.88,0.5,fill=SOFT,line=SOFTLINE,lw=1,rounded=True,radius=0.15,dash="dash")
text(s,8.72,4.32,3.6,0.5,[[("θt · trainable   ",{'bold':True,'color':ACCENT,'font':SERIF}),("8,396,806 params",{'color':ACCENT_D})]],size=13,anchor=MSO_ANCHOR.MIDDLE)
callout(s,0.9,5.25,11.5,1.35,[
    ("The subtle bug I avoided:  ",{'bold':True,'color':ACCENT}),
    ("if frozen BatchNorm layers keep updating their running stats, the representation drifts and training collapses. In both frozen conditions BN stays in eval mode. In Random-Full nothing is frozen, so BN updates normally — different by design, not by accident.",{}),
])
footer(s,"Transfer Learning vs. Random Init","04")

# ---------------- SLIDE 5 : SETUP ----------------
s=slide(); kicker(s,"Setup"); h2(s,"Same data, same recipe, three seeds")
cw=3.68; cx=0.9; cy=2.3; ch=4.0
card(s,cx,cy,cw,ch,"Data","Intel Scenes",[
    "6 classes: buildings, forest, glacier, mountain, sea, street",
    [("12,631",{'bold':True,'color':INK}),(" train / ",{}),("1,403",{'bold':True,'color':INK}),(" val / ",{}),("3,000",{'bold':True,'color':INK}),(" test (held out)",{})],
    "Resize 224, center-crop, ImageNet normalize",
])
card(s,cx+cw+0.19,cy,cw,ch,"Training","Identical for all",[
    [("SGD · lr ",{}),("1e-3",{'bold':True,'color':INK}),(" · momentum ",{}),("0.9",{'bold':True,'color':INK}),(" · wd ",{}),("1e-4",{'bold':True,'color':INK})],
    "Batch 32 · cross-entropy on θt",
    [("Early stop on val loss, patience ",{}),("5",{'bold':True,'color':INK})],
])
card(s,cx+2*(cw+0.19),cy,cw,ch,"Rigor","Reproducible",[
    [("3 seeds: ",{}),("42, 100, 2026",{'bold':True,'color':INK})],
    [("Test split touched ",{}),("only once",{'bold':True,'color':INK}),(", after model selection",{})],
    "Report mean ± std, not single runs",
])
footer(s,"Transfer Learning vs. Random Init","05")

# ---------------- SLIDE 6 : RESULTS (accuracy + cost) ----------------
s=slide(); kicker(s,"Results"); h2(s,"Pretrained wins — on accuracy and speed")
# left: accuracy bars
text(s,0.9,2.2,4.5,0.3,"TEST ACCURACY  ·  mean of 3 seeds",size=11,color=MUTED,bold=True)
bars=[("PT-F",93.3,ACCENT),("R-FL",84.7,SLATE6),("R-F",66.5,SLATE4)]
lx=0.9; lw=1.3; trackx=2.4; trackw=4.55; by=2.72; bh=0.6; vgap=0.28
for name,val,col in bars:
    text(s,lx,by,lw,bh,name,size=14,color=INK,bold=True,font=SERIF,align=PP_ALIGN.RIGHT,anchor=MSO_ANCHOR.MIDDLE)
    rect(s,trackx,by,trackw,bh,fill=FROZBG,rounded=True,radius=0.3)
    fw=trackw*val/100.0
    rect(s,trackx,by,fw,bh,fill=col,rounded=True,radius=0.3)
    text(s,trackx+fw-1.1,by,0.95,bh,f"{val}%",size=15,color=WHITE,bold=True,font=SERIF,align=PP_ALIGN.RIGHT,anchor=MSO_ANCHOR.MIDDLE)
    by+=bh+vgap
# right: convergence & cost panel
px,py,pw,ph=7.55,2.5,4.85,2.32
rect(s,px,py,pw,ph,fill=PANEL,line=RULE,lw=1,rounded=True,radius=0.05)
text(s,px+0.32,py+0.24,pw-0.6,0.3,"CONVERGENCE & COST",size=11,color=ACCENT,bold=True)
rows=[("Pretrained-Frozen","8.0 epochs · ~103s",ACCENT),
      ("Random-Frozen","8.7 epochs · ~113s",SLATE6),
      ("Random-Full","15.0 epochs · ~426s",SLATE6)]
ry=py+0.66
for nm,val,col in rows:
    text(s,px+0.32,ry,2.4,0.3,nm,size=13,color=INK,bold=True)
    text(s,px+2.55,ry,pw-2.85,0.3,val,size=13,color=col,align=PP_ALIGN.RIGHT)
    ry+=0.42
text(s,px+0.32,ry+0.04,pw-0.6,0.3,[[("Training from scratch pays ",{'color':MUTED}),("~4× the wall-clock.",{'bold':True,'color':INK})]],size=12.5)
# bottom callout
callout(s,0.9,5.4,11.5,1.15,[
    ("Capacity-matched gap: +26.7 points.  ",{'bold':True,'color':ACCENT}),
    ("PT-F vs R-F — identical trainable parameters and recipe, the only difference is pretrained vs random frozen features. That gap is pretraining, cleanly isolated. Unfreezing everything (R-FL) recovers accuracy but never reaches the pretrained ceiling.",{}),
])
footer(s,"Mean test accuracy over 3 seeds · std ≤ 4%","06")

IMGDIR=os.path.join(os.path.dirname(os.path.abspath(__file__)),"..","results")

# ---------------- SLIDE 7 : TRAINING DYNAMICS ----------------
s=slide(); kicker(s,"Results"); h2(s,"Pretrained converges instantly; scratch crawls",rule=False)
text(s,0.9,1.82,11.5,0.55,"Pretrained-frozen starts near its ceiling and flatlines in ~4 epochs; random-full climbs for 12–17 epochs; random-frozen is noisy and low.",size=14.5,color=MUTED,line_sp=1.3)
image_centered(s,os.path.join(IMGDIR,"learning_curves.png"),top=2.55,width=11.3)
footer(s,"Per-epoch validation metrics · 3 seeds","07")

# ---------------- SLIDE 8 : FEATURE SPACE ----------------
s=slide(); kicker(s,"Why It Works"); h2(s,"Pretrained features are already separable",rule=False)
text(s,0.9,1.82,11.5,0.55,"PCA of the 512-d penultimate features: pretrained-frozen splits the classes (overlaps match the confusion matrix); random-frozen stays entangled.",size=14.5,color=MUTED,line_sp=1.3)
image_centered(s,os.path.join(IMGDIR,"feature_pca.png"),top=2.5,width=9.9)
footer(s,"Test split · seed 42","08")

# ---------------- SLIDE 9 : ABLATION ----------------
s=slide(); kicker(s,"Ablation"); h2(s,"How much fine-tuning do you actually need?",rule=False)
from PIL import Image as _Img
_ab=os.path.join(IMGDIR,"ablation_curve.png"); _iw,_ih=_Img.open(_ab).size
_w=7.2; _h=_w*_ih/_iw
s.shapes.add_picture(_ab,Inches(0.7),Inches(2.55),Inches(_w),Inches(_h))
bullets(s,8.25,2.55,4.15,4.0,[
    [("A bare linear probe already hits ",{}),("91.3%",{'bold':True,'color':ACCENT}),(" — frozen ImageNet features are nearly linearly separable.",{})],
    [("Unfreezing ",{}),("layer4",{'bold':True,'color':INK}),(" adds ~1.6 pts: the boundary used in the main experiments.",{})],
    [("layer3, layer2, full fine-tune: ",{}),("flat at 92.8%",{'bold':True,'color':INK}),(" — 3.6× the params, no gain.",{})],
    [("Most of pretraining’s value needs ",{}),("almost no fine-tuning",{'bold':True,'color':INK}),(".",{})],
],size=14,gap=14)
footer(s,"Pretrained ResNet-18 · 40% train subset · seed 42","09")

# ---------------- SLIDE 10 : GABOR FIRST LAYER ----------------
s=slide(); kicker(s,"Extension"); h2(s,"How far up does the first layer reach?",rule=False)
_gb=os.path.join(IMGDIR,"gabor_curve.png"); _giw,_gih=_Img.open(_gb).size
_gw=7.4; _gh=_gw*_gih/_giw
s.shapes.add_picture(_gb,Inches(0.55),Inches(2.75),Inches(_gw),Inches(_gh))
bullets(s,8.15,2.7,4.25,4.0,[
    [("Swap ",{}),("conv1",{'bold':True,'color':INK}),(" for a fixed ",{}),("Gabor bank",{'bold':True,'color':INK}),(" (8 orientations × 4 freqs × 2 phases), freeze it, then unfreeze deeper stages one at a time.",{})],
    [("Gabor read-out: ",{}),("81.1%",{'bold':True,'color':ACCENT}),(" — about 10 pts below the learned first layer.",{})],
    [("Unfreezing just ",{}),("layer1",{'bold':True,'color':INK}),(" recovers ",{}),("+5 pts to 86.1%",{'bold':True,'color':INK}),("; layer2–4 add almost nothing.",{})],
    [("First-layer influence reaches ",{}),("~one stage up",{'bold':True,'color':ACCENT}),(". A ~7-pt gap to learned features never closes, even fully trainable.",{})],
],size=13.5,gap=13)
footer(s,"Gabor vs. learned conv1 · full data · seed 42","10")

# ---------------- SLIDE 11 : CONFUSION ----------------
s=slide(); kicker(s,"Error Analysis"); h2(s,"The mistakes are the dataset’s, not the method’s")
labels=["B","F","G","M","Se","St"]
cm=[[408,0,0,0,1,28],[0,471,0,1,0,2],[1,4,458,71,16,3],[1,2,28,484,9,1],[2,0,2,3,503,0],[26,0,0,2,0,473]]
hot={(0,5),(2,3),(5,0)}
gx=1.0; gy=2.75; cell=0.52; hdr=0.42
text(s,gx,gy-0.42,5,0.3,"Pretrained-Frozen · test confusion (seed 100)",size=13,color=INK,bold=True,font=SERIF)
for j,lb in enumerate(labels):
    text(s,gx+hdr+j*cell,gy-0.02,cell,0.3,lb,size=11,color=MUTED,bold=True,align=PP_ALIGN.CENTER)
for i,row in enumerate(cm):
    ry=gy+hdr+i*cell
    text(s,gx-0.05,ry,hdr,cell,labels[i],size=11,color=MUTED,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    for j,v in enumerate(row):
        cxp=gx+hdr+j*cell
        if i==j:
            rect(s,cxp,ry,cell,cell,fill=ACCENT,line=WHITE,lw=1)
            text(s,cxp,ry,cell,cell,str(v),size=11,color=WHITE,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        elif (i,j) in hot:
            rect(s,cxp,ry,cell,cell,fill=HOTBG,line=WHITE,lw=1)
            text(s,cxp,ry,cell,cell,str(v),size=11,color=ACCENT,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        else:
            rect(s,cxp,ry,cell,cell,fill=WHITE,line=FROZBG,lw=1)
            text(s,cxp,ry,cell,cell,str(v),size=10.5,color=FAINT,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
bullets(s,7.5,2.95,4.9,3.6,[
    [("Glacier ↔ Mountain",{'bold':True,'color':INK}),(" (71 errors): shared snow, ridges, rock texture — genuinely ambiguous.",{})],
    [("Buildings ↔ Street",{'bold':True,'color':INK}),(" (26–28): streets contain buildings and vice-versa.",{})],
    [("The ",{}),("same",{'bold':True,'color':INK}),(" confusions appear in all three conditions — just far larger for random init.",{})],
    [("Errors track ",{}),("visual overlap in the data",{'bold':True,'color':INK}),(", not a broken training setup.",{})],
],size=15,gap=13)
footer(s,"B buildings · F forest · G glacier · M mountain · Se sea · St street","11")

# ---------------- SLIDE 12 : TAKEAWAYS ----------------
s=slide(); kicker(s,"Conclusion"); h2(s,"What the control actually proved")
takes=[
    ("01","Pretraining is the feature quality, not the parameter count",
     [("Capacity-matched, pretrained features add ",{}),("+26.7 points",{'bold':True,'color':ACCENT}),(" over random frozen features. The win is the representations themselves.",{})]),
    ("02","Capacity helps, but can’t close the gap",
     [("Unfreezing everything (Random-Full) recovers to 84.7% — better than random-frozen, yet still ~9 points short of pretraining, at ~4× the compute.",{})]),
    ("03","For small data + tight budgets, transfer learning wins outright",
     [("Best accuracy and fastest convergence. The random-frozen control is what lets me say that with confidence.",{})]),
]
ty=2.4
for num,title,body in takes:
    text(s,0.9,ty,0.9,0.8,num,size=30,color=SOFTLINE,bold=True,font=SERIF)
    text(s,1.85,ty,10.5,0.4,title,size=19,color=INK,bold=True,font=SERIF)
    text(s,1.85,ty+0.48,10.5,0.7,[body],size=14.5,color=SLATE6,line_sp=1.3)
    ty+=1.45
footer(s,"Thank you · Questions welcome","12")

prs.save(OUT)
print("saved",OUT)
